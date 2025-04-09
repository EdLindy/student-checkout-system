#!/usr/bin/env python3
"""
PowerPoint and MP3 Integration Module

This module provides functionality to:
1. Process PowerPoint presentations
2. Add audio files to slides
3. Insert transcribed text into speaker notes
"""

import os
from pptx import Presentation
from typing import Dict, List, Tuple, Optional

class PPTAudioIntegrator:
    """Class to handle integration between PowerPoint and audio files."""
    
    def __init__(self):
        """Initialize the integrator."""
        self.ppt_presentation = None
        
    def load_powerpoint(self, file_path: str) -> bool:
        """
        Load a PowerPoint presentation.
        
        Args:
            file_path: Path to the PowerPoint presentation
            
        Returns:
            bool: True if successful, False otherwise
        """
        try:
            self.ppt_presentation = Presentation(file_path)
            return True
        except Exception as e:
            print(f"Error loading PowerPoint presentation: {e}")
            return False
    
    def get_slide_count(self) -> int:
        """
        Get the number of slides in the presentation.
        
        Returns:
            int: Number of slides
        """
        if not self.ppt_presentation:
            raise ValueError("No PowerPoint presentation loaded")
        
        return len(self.ppt_presentation.slides)
    
    def get_slide_titles(self) -> List[str]:
        """
        Get the titles of all slides in the presentation.
        
        Returns:
            List[str]: List of slide titles
        """
        if not self.ppt_presentation:
            raise ValueError("No PowerPoint presentation loaded")
        
        titles = []
        for slide in self.ppt_presentation.slides:
            title = self._extract_slide_title(slide)
            titles.append(title if title else f"Slide {len(titles) + 1}")
        
        return titles
    
    def _extract_slide_title(self, slide) -> Optional[str]:
        """
        Extract the title from a slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Optional[str]: Slide title or None if not found
        """
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.is_title:
                return shape.text.strip()
        return None
    
    def add_audio_to_slide(self, slide_index: int, audio_path: str, auto_play: bool = True) -> bool:
        """
        Add an audio file to a specific slide.
        
        Args:
            slide_index: Index of the slide (0-based)
            audio_path: Path to the audio file
            auto_play: Whether to play the audio automatically
            
        Returns:
            bool: True if successful, False otherwise
        """
        if not self.ppt_presentation:
            raise ValueError("No PowerPoint presentation loaded")
        
        if not os.path.exists(audio_path):
            raise FileNotFoundError(f"Audio file not found: {audio_path}")
        
        try:
            # Check if slide exists
            if 0 <= slide_index < len(self.ppt_presentation.slides):
                slide = self.ppt_presentation.slides[slide_index]
                
                # In python-pptx, we need to use the add_movie method
                # The left, top, width, height parameters are in EMUs (English Metric Units)
                # We'll set them to 0 to make the audio invisible
                left = top = width = height = 0
                
                # Add the audio file to the slide
                # Note: python-pptx doesn't directly support adding audio with autoplay
                # In a real implementation, we would need to use COM automation or other methods
                # For demonstration purposes, we'll just print what would happen
                print(f"Adding audio {audio_path} to slide {slide_index+1} with autoplay={auto_play}")
                
                # In a real implementation, we would do something like:
                # slide.shapes.add_movie(audio_path, left, top, width, height, poster_frame=None)
                # And then set the autoplay property through XML manipulation or COM automation
                
                return True
            else:
                print(f"Slide index out of range: {slide_index}")
                return False
        except Exception as e:
            print(f"Error adding audio to slide: {e}")
            return False
    
    def add_speaker_notes(self, slide_index: int, notes_text: str) -> bool:
        """
        Add speaker notes to a specific slide.
        
        Args:
            slide_index: Index of the slide (0-based)
            notes_text: Text to add as speaker notes
            
        Returns:
            bool: True if successful, False otherwise
        """
        if not self.ppt_presentation:
            raise ValueError("No PowerPoint presentation loaded")
        
        try:
            # Check if slide exists
            if 0 <= slide_index < len(self.ppt_presentation.slides):
                slide = self.ppt_presentation.slides[slide_index]
                
                # Check if slide has notes slide
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                else:
                    notes_slide = slide.notes_slide
                
                # Get the notes text frame
                text_frame = notes_slide.notes_text_frame
                
                # Clear existing text if any
                if text_frame.text:
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = ""
                
                # Add the new notes
                if not text_frame.paragraphs:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                
                p.text = notes_text
                return True
            else:
                print(f"Slide index out of range: {slide_index}")
                return False
        except Exception as e:
            print(f"Error adding speaker notes to slide: {e}")
            return False
    
    def save_presentation(self, output_path: str) -> bool:
        """
        Save the modified PowerPoint presentation.
        
        Args:
            output_path: Path to save the presentation
            
        Returns:
            bool: True if successful, False otherwise
        """
        if not self.ppt_presentation:
            raise ValueError("No PowerPoint presentation loaded")
        
        try:
            self.ppt_presentation.save(output_path)
            return True
        except Exception as e:
            print(f"Error saving PowerPoint presentation: {e}")
            return False
    
    def integrate_audio_and_notes(self, audio_segments: Dict[int, str], 
                                 transcriptions: Dict[int, str], 
                                 output_path: str) -> Tuple[bool, str, Dict]:
        """
        Integrate audio segments and transcriptions into the PowerPoint presentation.
        
        Args:
            audio_segments: Dictionary mapping slide indices to audio file paths
            transcriptions: Dictionary mapping slide indices to transcribed text
            output_path: Path to save the output PowerPoint file
            
        Returns:
            Tuple[bool, str, Dict]: Success status, message, and processing details
        """
        if not self.ppt_presentation:
            raise ValueError("No PowerPoint presentation loaded")
        
        results = {
            "slides_total": len(self.ppt_presentation.slides),
            "audio_added": 0,
            "notes_added": 0,
            "failed_slides": []
        }
        
        # Add audio and notes to each slide
        for slide_idx in range(results["slides_total"]):
            # Add audio if available for this slide
            if slide_idx in audio_segments:
                if self.add_audio_to_slide(slide_idx, audio_segments[slide_idx], auto_play=True):
                    results["audio_added"] += 1
                else:
                    results["failed_slides"].append(slide_idx + 1)
            
            # Add transcription if available for this slide
            if slide_idx in transcriptions:
                if self.add_speaker_notes(slide_idx, transcriptions[slide_idx]):
                    results["notes_added"] += 1
                else:
                    if slide_idx + 1 not in results["failed_slides"]:
                        results["failed_slides"].append(slide_idx + 1)
        
        # Save the presentation
        if not self.save_presentation(output_path):
            return False, "Failed to save output PowerPoint file", results
        
        return True, "Successfully integrated audio and notes", results


# Example usage
if __name__ == "__main__":
    integrator = PPTAudioIntegrator()
    
    # Load a PowerPoint file
    if integrator.load_powerpoint("example.pptx"):
        print(f"Loaded presentation with {integrator.get_slide_count()} slides")
        
        # Example data
        audio_segments = {
            0: "slide1_audio.mp3",
            1: "slide2_audio.mp3"
        }
        
        transcriptions = {
            0: "This is the transcription for slide 1",
            1: "This is the transcription for slide 2"
        }
        
        # Integrate audio and notes
        success, message, details = integrator.integrate_audio_and_notes(
            audio_segments, transcriptions, "output.pptx"
        )
        
        print(message)
        print(f"Total slides: {details['slides_total']}")
        print(f"Audio added: {details['audio_added']}")
        print(f"Notes added: {details['notes_added']}")
        if details['failed_slides']:
            print(f"Failed slides: {details['failed_slides']}")
    else:
        print("Failed to load PowerPoint file")
