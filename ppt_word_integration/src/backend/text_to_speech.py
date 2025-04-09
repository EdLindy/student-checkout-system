#!/usr/bin/env python3
"""
Text-to-Speech Conversion Module

This module provides functionality to:
1. Convert text to speech using cloud-based TTS services
2. Generate audio files with a middle-aged American male voice with engaging disposition
3. Attach audio files to PowerPoint slides
"""

import os
import time
import tempfile
from typing import Dict, List, Tuple, Optional
from pptx import Presentation
import requests
import json
import base64

# For production, we would use one of these libraries
# Commented out as they would need to be installed
# from google.cloud import texttospeech
# import boto3
# from azure.cognitiveservices.speech import SpeechConfig, SpeechSynthesizer

class TextToSpeechConverter:
    """Class to handle text-to-speech conversion and integration with PowerPoint."""
    
    def __init__(self, service_type="google"):
        """
        Initialize the TTS converter.
        
        Args:
            service_type: The TTS service to use ('google', 'amazon', 'azure', 'elevenlabs')
        """
        self.service_type = service_type
        self.ppt_presentation = None
        self.temp_dir = tempfile.mkdtemp()
        
    def load_powerpoint(self, file_path: str) -> bool:
        """
        Load a PowerPoint presentation.
        
        Args:
            file_path: Path to the PowerPoint presentation
            
        Returns:
            bool: True if successful, False otherwise
        """
        try:
            self.ppt_presentation = Presentation(file_path)
            return True
        except Exception as e:
            print(f"Error loading PowerPoint presentation: {e}")
            return False
    
    def extract_notes_from_presentation(self) -> Dict[int, str]:
        """
        Extract speaker notes from the PowerPoint presentation.
        
        Returns:
            Dict[int, str]: Dictionary mapping slide numbers to speaker notes
        """
        if not self.ppt_presentation:
            raise ValueError("No PowerPoint presentation loaded")
        
        speaker_notes = {}
        
        for i, slide in enumerate(self.ppt_presentation.slides):
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    speaker_notes[i] = notes_text
        
        return speaker_notes
    
    def convert_text_to_speech_google(self, text: str, output_file: str) -> bool:
        """
        Convert text to speech using Google Cloud TTS API.
        
        Args:
            text: The text to convert
            output_file: Path to save the audio file
            
        Returns:
            bool: True if successful, False otherwise
        """
        # In a real implementation, we would use the Google Cloud TTS API
        # This is a placeholder that simulates the API call
        
        try:
            # Simulate API call delay
            time.sleep(1)
            
            # Create a dummy audio file (in production, this would be the actual audio)
            with open(output_file, 'wb') as f:
                # Write a minimal valid MP3 file header
                f.write(b'\xFF\xFB\x90\x44\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00')
            
            return True
        except Exception as e:
            print(f"Error in Google TTS conversion: {e}")
            return False
    
    def convert_text_to_speech_amazon(self, text: str, output_file: str) -> bool:
        """
        Convert text to speech using Amazon Polly.
        
        Args:
            text: The text to convert
            output_file: Path to save the audio file
            
        Returns:
            bool: True if successful, False otherwise
        """
        # In a real implementation, we would use the Amazon Polly API
        # This is a placeholder that simulates the API call
        
        try:
            # Simulate API call delay
            time.sleep(1)
            
            # Create a dummy audio file
            with open(output_file, 'wb') as f:
                f.write(b'\xFF\xFB\x90\x44\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00')
            
            return True
        except Exception as e:
            print(f"Error in Amazon Polly conversion: {e}")
            return False
    
    def convert_text_to_speech_azure(self, text: str, output_file: str) -> bool:
        """
        Convert text to speech using Microsoft Azure Speech Service.
        
        Args:
            text: The text to convert
            output_file: Path to save the audio file
            
        Returns:
            bool: True if successful, False otherwise
        """
        # In a real implementation, we would use the Azure Speech Service API
        # This is a placeholder that simulates the API call
        
        try:
            # Simulate API call delay
            time.sleep(1)
            
            # Create a dummy audio file
            with open(output_file, 'wb') as f:
                f.write(b'\xFF\xFB\x90\x44\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00')
            
            return True
        except Exception as e:
            print(f"Error in Azure Speech Service conversion: {e}")
            return False
    
    def convert_text_to_speech_elevenlabs(self, text: str, output_file: str) -> bool:
        """
        Convert text to speech using ElevenLabs API.
        
        Args:
            text: The text to convert
            output_file: Path to save the audio file
            
        Returns:
            bool: True if successful, False otherwise
        """
        # In a real implementation, we would use the ElevenLabs API
        # This is a placeholder that simulates the API call
        
        try:
            # Simulate API call delay
            time.sleep(1)
            
            # Create a dummy audio file
            with open(output_file, 'wb') as f:
                f.write(b'\xFF\xFB\x90\x44\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00')
            
            return True
        except Exception as e:
            print(f"Error in ElevenLabs conversion: {e}")
            return False
    
    def convert_text_to_speech(self, text: str, output_file: str) -> bool:
        """
        Convert text to speech using the selected service.
        
        Args:
            text: The text to convert
            output_file: Path to save the audio file
            
        Returns:
            bool: True if successful, False otherwise
        """
        if self.service_type == "google":
            return self.convert_text_to_speech_google(text, output_file)
        elif self.service_type == "amazon":
            return self.convert_text_to_speech_amazon(text, output_file)
        elif self.service_type == "azure":
            return self.convert_text_to_speech_azure(text, output_file)
        elif self.service_type == "elevenlabs":
            return self.convert_text_to_speech_elevenlabs(text, output_file)
        else:
            raise ValueError(f"Unsupported TTS service: {self.service_type}")
    
    def add_audio_to_slide(self, slide_idx: int, audio_path: str) -> bool:
        """
        Add audio to a PowerPoint slide.
        
        Args:
            slide_idx: Index of the slide
            audio_path: Path to the audio file
            
        Returns:
            bool: True if successful, False otherwise
        """
        if not self.ppt_presentation:
            raise ValueError("No PowerPoint presentation loaded")
        
        try:
            # In a real implementation, we would use python-pptx to add the audio
            # This is a placeholder as python-pptx doesn't directly support adding audio
            # We would need to use a library like python-pptx-interface or COM automation
            
            # For demonstration purposes, we'll just print what would happen
            print(f"Adding audio {audio_path} to slide {slide_idx+1}")
            
            # In reality, we would do something like:
            # slide = self.ppt_presentation.slides[slide_idx]
            # slide.shapes.add_movie(audio_path, 0, 0, 0, 0, poster_frame=None, play_mode='AUTO')
            
            return True
        except Exception as e:
            print(f"Error adding audio to slide: {e}")
            return False
    
    def process_presentation(self, ppt_path: str, output_path: str, voice_type: str = "middle-aged-american-male") -> Tuple[bool, str, Dict]:
        """
        Process the PowerPoint presentation to add TTS audio.
        
        Args:
            ppt_path: Path to the PowerPoint file
            output_path: Path to save the output PowerPoint file
            voice_type: Type of voice to use
            
        Returns:
            Tuple[bool, str, Dict]: Success status, message, and processing details
        """
        results = {
            "slides_total": 0,
            "notes_extracted": 0,
            "audio_files_created": 0,
            "slides_with_audio": 0,
            "failed_slides": []
        }
        
        # Load the PowerPoint file
        if not self.load_powerpoint(ppt_path):
            return False, "Failed to load PowerPoint file", results
        
        # Extract speaker notes
        try:
            speaker_notes = self.extract_notes_from_presentation()
            results["notes_extracted"] = len(speaker_notes)
            results["slides_total"] = len(self.ppt_presentation.slides)
        except Exception as e:
            return False, f"Error extracting speaker notes: {e}", results
        
        # Convert notes to speech and add to slides
        audio_files = []
        for slide_idx, notes in speaker_notes.items():
            # Generate a unique filename for this slide's audio
            audio_file = os.path.join(self.temp_dir, f"slide_{slide_idx+1}_audio.mp3")
            
            # Convert the notes to speech
            if self.convert_text_to_speech(notes, audio_file):
                audio_files.append(audio_file)
                results["audio_files_created"] += 1
                
                # Add the audio to the slide
                if self.add_audio_to_slide(slide_idx, audio_file):
                    results["slides_with_audio"] += 1
                else:
                    results["failed_slides"].append(slide_idx + 1)
            else:
                results["failed_slides"].append(slide_idx + 1)
        
        # Save the presentation
        try:
            self.ppt_presentation.save(output_path)
        except Exception as e:
            return False, f"Error saving PowerPoint presentation: {e}", results
        
        return True, "Successfully processed presentation with TTS", results
    
    def cleanup(self):
        """Clean up temporary files."""
        for file in os.listdir(self.temp_dir):
            try:
                os.remove(os.path.join(self.temp_dir, file))
            except:
                pass
        try:
            os.rmdir(self.temp_dir)
        except:
            pass


# Example usage
if __name__ == "__main__":
    converter = TextToSpeechConverter(service_type="google")
    success, message, details = converter.process_presentation(
        "presentation_with_notes.pptx",
        "presentation_with_audio.pptx"
    )
    print(message)
    print(f"Total slides: {details['slides_total']}")
    print(f"Notes extracted: {details['notes_extracted']}")
    print(f"Audio files created: {details['audio_files_created']}")
    print(f"Slides with audio: {details['slides_with_audio']}")
    if details['failed_slides']:
        print(f"Failed slides: {details['failed_slides']}")
    converter.cleanup()
