#!/usr/bin/env python3
"""
PowerPoint and Word Integration Module

This module provides functionality to:
1. Extract speaker notes from Word documents
2. Insert these notes into PowerPoint presentations
"""

import re
import os
from pptx import Presentation
from docx import Document
from typing import Dict, Tuple, List, Optional

class PPTWordIntegrator:
    """Class to handle integration between PowerPoint and Word files."""
    
    def __init__(self):
        """Initialize the integrator."""
        self.word_doc = None
        self.ppt_presentation = None
        
    def load_word_document(self, file_path: str) -> bool:
        """
        Load a Word document.
        
        Args:
            file_path: Path to the Word document
            
        Returns:
            bool: True if successful, False otherwise
        """
        try:
            self.word_doc = Document(file_path)
            return True
        except Exception as e:
            print(f"Error loading Word document: {e}")
            return False
    
    def load_powerpoint(self, file_path: str) -> bool:
        """
        Load a PowerPoint presentation.
        
        Args:
            file_path: Path to the PowerPoint presentation
            
        Returns:
            bool: True if successful, False otherwise
        """
        try:
            self.ppt_presentation = Presentation(file_path)
            return True
        except Exception as e:
            print(f"Error loading PowerPoint presentation: {e}")
            return False
    
    def extract_speaker_notes(self) -> Dict[int, str]:
        """
        Extract speaker notes from the Word document.
        
        Returns:
            Dict[int, str]: Dictionary mapping slide numbers to speaker notes
        """
        if not self.word_doc:
            raise ValueError("No Word document loaded")
        
        speaker_notes = {}
        current_slide = None
        current_text = []
        
        # Pattern to match "Slide X" headers
        slide_pattern = re.compile(r'^Slide\s+(\d+)', re.IGNORECASE)
        
        for paragraph in self.word_doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
                
            # Check if this paragraph is a slide header
            match = slide_pattern.match(text)
            if match:
                # If we were collecting text for a previous slide, save it
                if current_slide is not None and current_text:
                    speaker_notes[current_slide] = '\n'.join(current_text)
                    current_text = []
                
                # Start collecting text for this new slide
                current_slide = int(match.group(1))
            elif current_slide is not None:
                # Add this paragraph to the current slide's notes
                current_text.append(text)
        
        # Don't forget to save the last slide's notes
        if current_slide is not None and current_text:
            speaker_notes[current_slide] = '\n'.join(current_text)
        
        return speaker_notes
    
    def insert_speaker_notes(self, speaker_notes: Dict[int, str]) -> Tuple[int, List[int]]:
        """
        Insert speaker notes into the PowerPoint presentation.
        
        Args:
            speaker_notes: Dictionary mapping slide numbers to speaker notes
            
        Returns:
            Tuple[int, List[int]]: Count of updated slides and list of slide numbers that couldn't be updated
        """
        if not self.ppt_presentation:
            raise ValueError("No PowerPoint presentation loaded")
        
        updated_count = 0
        failed_slides = []
        
        # PowerPoint slides are 0-indexed, but users typically refer to them as 1-indexed
        for slide_num, notes in speaker_notes.items():
            # Convert to 0-indexed
            idx = slide_num - 1
            
            # Check if this slide exists
            if 0 <= idx < len(self.ppt_presentation.slides):
                slide = self.ppt_presentation.slides[idx]
                
                # Check if slide has notes slide
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                else:
                    notes_slide = slide.notes_slide
                
                # Get the notes text frame
                text_frame = notes_slide.notes_text_frame
                
                # Clear existing text if any
                if text_frame.text:
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = ""
                
                # Add the new notes
                if not text_frame.paragraphs:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                
                p.text = notes
                updated_count += 1
            else:
                failed_slides.append(slide_num)
        
        return updated_count, failed_slides
    
    def save_presentation(self, output_path: str) -> bool:
        """
        Save the modified PowerPoint presentation.
        
        Args:
            output_path: Path to save the presentation
            
        Returns:
            bool: True if successful, False otherwise
        """
        if not self.ppt_presentation:
            raise ValueError("No PowerPoint presentation loaded")
        
        try:
            self.ppt_presentation.save(output_path)
            return True
        except Exception as e:
            print(f"Error saving PowerPoint presentation: {e}")
            return False
    
    def process_files(self, ppt_path: str, word_path: str, output_path: str) -> Tuple[bool, str, Dict]:
        """
        Process the PowerPoint and Word files.
        
        Args:
            ppt_path: Path to the PowerPoint file
            word_path: Path to the Word file
            output_path: Path to save the output PowerPoint file
            
        Returns:
            Tuple[bool, str, Dict]: Success status, message, and processing details
        """
        results = {
            "slides_total": 0,
            "notes_extracted": 0,
            "slides_updated": 0,
            "failed_slides": []
        }
        
        # Load the files
        if not self.load_powerpoint(ppt_path):
            return False, "Failed to load PowerPoint file", results
        
        if not self.load_word_document(word_path):
            return False, "Failed to load Word document", results
        
        # Extract speaker notes
        try:
            speaker_notes = self.extract_speaker_notes()
            results["notes_extracted"] = len(speaker_notes)
        except Exception as e:
            return False, f"Error extracting speaker notes: {e}", results
        
        # Insert speaker notes
        try:
            results["slides_total"] = len(self.ppt_presentation.slides)
            updated_count, failed_slides = self.insert_speaker_notes(speaker_notes)
            results["slides_updated"] = updated_count
            results["failed_slides"] = failed_slides
        except Exception as e:
            return False, f"Error inserting speaker notes: {e}", results
        
        # Save the presentation
        if not self.save_presentation(output_path):
            return False, "Failed to save output PowerPoint file", results
        
        return True, "Successfully processed files", results


# Example usage
if __name__ == "__main__":
    integrator = PPTWordIntegrator()
    success, message, details = integrator.process_files(
        "example.pptx", 
        "speaker_notes.docx", 
        "output.pptx"
    )
    print(message)
    print(f"Total slides: {details['slides_total']}")
    print(f"Notes extracted: {details['notes_extracted']}")
    print(f"Slides updated: {details['slides_updated']}")
    if details['failed_slides']:
        print(f"Failed slides: {details['failed_slides']}")
